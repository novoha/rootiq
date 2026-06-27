# Generates RootIQ-Presentation.pptx (editable PowerPoint) from project content.
#   python make_pptx.py
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SHOT = lambda n: os.path.abspath(f"working_dir/screenshots/{n}")
GH = "https://github.com/novoha/rootiq"
APP = "https://rootiq.streamlit.app/"

NAVY = RGBColor(0x0D, 0x1B, 0x2A)
NAVY2 = RGBColor(0x08, 0x13, 0x20)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEALL = RGBColor(0x5E, 0xEA, 0xD4)
INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x64, 0x74, 0x8B)
MUTED2 = RGBColor(0x94, 0xA3, 0xB8)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
SOFT = RGBColor(0xF7, 0xF9, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OK = RGBColor(0x16, 0x65, 0x34)
OKBG = RGBColor(0xDC, 0xFC, 0xE7)
WARN = RGBColor(0x92, 0x40, 0x0E)
WARNBG = RGBColor(0xFE, 0xF3, 0xC7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill, line=None, rounded=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def para(tf, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, after=6, first=False, line=1.12):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(after); p.space_before = Pt(0); p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, bold, color)]
    for txt, b, c in runs:
        r = p.add_run(); r.text = txt
        f = r.font; f.size = Pt(size); f.bold = b; f.color.rgb = c; f.name = "Segoe UI"
    return p


def bullet(tf, txt, lead=None, size=17, after=9):
    p = tf.add_paragraph(); p.space_after = Pt(after); p.line_spacing = 1.12
    sq = p.add_run(); sq.text = "▪  "; sq.font.size = Pt(size); sq.font.color.rgb = TEAL; sq.font.bold = True; sq.font.name = "Segoe UI"
    if lead:
        r = p.add_run(); r.text = lead; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = INK; r.font.name = "Segoe UI"
    r = p.add_run(); r.text = txt; r.font.size = Pt(size); r.font.color.rgb = INK; r.font.name = "Segoe UI"
    return p


def link(tf, label, url, color=WHITE, size=15, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    r = p.add_run(); r.text = label; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = color; r.font.name = "Segoe UI"; r.hyperlink.address = url
    return p


def header(s, kicker, title):
    tf = tb(s, 0.75, 0.5, 11.8, 1.5)
    para(tf, kicker.upper(), 13, TEAL, True, after=5, first=True)
    para(tf, title, 30, NAVY, True, after=0, line=1.05)


def footer(s, n):
    rect(s, 0.75, 6.92, 11.83, 0.012, LINE)
    t1 = tb(s, 0.75, 7.0, 7, 0.35)
    para(t1, [("RootIQ, ", False, MUTED2), ("PLC Error Intelligence", True, TEAL)], 10, after=0, first=True)
    t2 = tb(s, 8.0, 7.0, 4.58, 0.35)
    para(t2, f"rootiq.streamlit.app    {n} / 15", 10, MUTED2, align=PP_ALIGN.RIGHT, after=0, first=True)


def card(s, x, y, w, h, title, desc):
    rect(s, x, y, w, h, WHITE, LINE, rounded=True)
    tf = tb(s, x + 0.28, y + 0.26, w - 0.56, h - 0.5)
    para(tf, title, 17, NAVY, True, after=5, first=True)
    para(tf, desc, 12.5, MUTED, after=0, line=1.16)


def img_fit(s, path, x, y, maxw, maxh):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = maxw; h = w / ar
    if h > maxh:
        h = maxh; w = h * ar
    rect(s, x - 0.04, y - 0.04, w + 0.08, h + 0.08, LINE, rounded=True)
    s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def links_line(s, x, y, color=WHITE):
    tf = tb(s, x, y, 11, 0.45)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "GitHub  "; r.font.size = Pt(15); r.font.color.rgb = color; r.font.bold = True; r.font.name = "Segoe UI"
    r = p.add_run(); r.text = "github.com/novoha/rootiq"; r.font.size = Pt(15); r.font.color.rgb = TEALL; r.font.bold = True; r.font.name = "Segoe UI"; r.hyperlink.address = GH
    r = p.add_run(); r.text = "      Live app  "; r.font.size = Pt(15); r.font.color.rgb = color; r.font.bold = True; r.font.name = "Segoe UI"
    r = p.add_run(); r.text = "rootiq.streamlit.app"; r.font.size = Pt(15); r.font.color.rgb = TEALL; r.font.bold = True; r.font.name = "Segoe UI"; r.hyperlink.address = APP


# ---------- 1 COVER ----------
s = slide(); rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 0, 0.22, 7.5, TEAL)
tf = tb(s, 0.95, 1.35, 11, 0.5)
para(tf, "INDUSTRIAL AI, 2026", 14, TEALL, True, after=0, first=True)
tf = tb(s, 0.9, 1.85, 11.5, 1.5)
para(tf, [("Root", True, WHITE), ("IQ", True, TEALL)], 66, after=0, first=True)
tf = tb(s, 0.95, 3.35, 11.4, 0.6)
para(tf, "PLC Error Intelligence, a secure agent harness", 24, RGBColor(0xE2, 0xE8, 0xF0), True, after=0, first=True)
tf = tb(s, 0.95, 4.05, 10.8, 1.1)
para(tf, "RootIQ reads a controller fault log, reuses past fixes, mines the IQAN community forum, and synthesises a cited solution. Every action is screened by a security interlock before it runs.", 15, MUTED2, after=0, first=True, line=1.3)
tf = tb(s, 0.95, 5.35, 11, 0.4)
para(tf, "Victor Otu Hayford", 18, WHITE, True, after=0, first=True)
links_line(s, 0.95, 5.9)

# ---------- 2 PROBLEM ----------
s = slide(); header(s, "The problem", "A cryptic fault, and the fix is buried in a forum")
rect(s, 0.75, 2.25, 5.7, 2.3, NAVY, rounded=True)
tf = tb(s, 1.05, 2.5, 5.2, 1.9)
para(tf, "# IQAN controller system log", 13, MUTED2, first=True, after=8)
for code, val in [("XC44 ;", "No contact"), ("COUT-Suction hose 1 ;", "Open load"),
                  ("VIN-Joystick (C1:26) ;", "Low error"), ("MC43FS[0] ;", "No contact")]:
    para(tf, [(code + " ", False, RGBColor(0xE2, 0xE8, 0xF0)), (val, True, RGBColor(0xFB, 0xBF, 0x24))], 14, after=4)
tf = tb(s, 6.9, 2.35, 5.6, 2.3)
bullet(tf, "The answer lives across forum threads and comments, slow to find by hand", size=16)
bullet(tf, "The same faults recur, yet nothing is captured or reused", size=16)
bullet(tf, "Sensitive sites need it offline, with no cloud and no API key", size=16)
rect(s, 0.75, 5.1, 11.83, 0.9, SOFT, rounded=True)
rect(s, 0.75, 5.1, 0.07, 0.9, TEAL)
tf = tb(s, 1.05, 5.32, 11.2, 0.6, MSO_ANCHOR.MIDDLE)
para(tf, "RootIQ turns that moment into a guided, cited fix, and remembers it for next time.", 15, RGBColor(0x47, 0x55, 0x69), first=True, after=0)
footer(s, 2)

# ---------- 3 WHAT IT IS ----------
s = slide(); header(s, "What RootIQ is", "A reusable, secure agent harness for fault diagnosis")
caps = [("Intake", "Upload a PDF or image, type a code, or load a whole CSV export"),
        ("Reuse", "Instant history cache hit, with no LLM and no network"),
        ("Retrieve", "Offline forum index, then scrape threads and comments"),
        ("Synthesise", "The LLM reasons across sources into a cited JSON fix"),
        ("Remember", "Auto-saved, then a knowledge map of related faults"),
        ("Guard", "security_check() screens every tool call first")]
cx, cy, cw, ch = 0.75, 2.05, 3.83, 1.55
for i, (t, d) in enumerate(caps):
    x = cx + (i % 3) * (cw + 0.18); y = cy + (i // 3) * (ch + 0.2)
    card(s, x, y, cw, ch, t, d)
tf = tb(s, 0.75, 5.55, 12, 0.5)
para(tf, [("Offline, local Ollama       ", True, TEAL), ("Online, any OpenAI-compatible API       ", True, TEAL), ("One sidebar toggle", True, TEAL)], 15, first=True, after=0)
footer(s, 3)

# ---------- 4 WHY IT MATTERS ----------
s = slide(); header(s, "Why it matters", "A working bridge between PLC programming and software")
rect(s, 0.75, 2.15, 5.8, 2.6, WHITE, LINE, rounded=True)
tf = tb(s, 1.05, 2.4, 5.25, 2.2)
para(tf, "Known faults resolve instantly", 17, NAVY, True, first=True, after=6)
para(tf, "The history cache turns a recurring fault into a zero-effort answer, with no LLM and no network. The same faults come back, and now they are solved on sight. This saves the most time.", 14, MUTED, after=0, line=1.25)
rect(s, 6.78, 2.15, 5.8, 2.6, WHITE, LINE, rounded=True)
tf = tb(s, 7.08, 2.4, 5.25, 2.2)
para(tf, "New faults are reasoned and searched", 17, NAVY, True, first=True, after=6)
para(tf, "Unseen faults are searched against the forum index, scraped with their community comments, and reasoned into a cited fix. Then they are remembered, so next time they are known too.", 14, MUTED, after=0, line=1.25)
rect(s, 0.75, 5.05, 11.83, 0.95, SOFT, rounded=True); rect(s, 0.75, 5.05, 0.07, 0.95, TEAL)
tf = tb(s, 1.05, 5.25, 11.2, 0.6, MSO_ANCHOR.MIDDLE)
para(tf, "The knowledge compounds. Every diagnosis grows the cache and the map, so the tool gets faster and more useful the more it runs.", 15, RGBColor(0x47, 0x55, 0x69), first=True, after=0)
footer(s, 4)

# ---------- 5 PIPELINE ----------
s = slide(); header(s, "How it works", "An orchestrated pipeline. The harness routes; the LLM only synthesises")
steps = [("Intake", "upload, type, CSV"), ("History", "instant reuse"), ("Search", "offline index"),
         ("Scrape", "threads, comments"), ("Synthesise", "cited JSON"), ("Remember", "map, history")]
sx, sw = 0.75, 1.78
for i, (t, d) in enumerate(steps):
    x = sx + i * (sw + 0.22)
    rect(s, x, 2.35, sw, 1.35, WHITE, LINE, rounded=True)
    tf = tb(s, x + 0.1, 2.6, sw - 0.2, 1.0)
    para(tf, t, 15, NAVY, True, align=PP_ALIGN.CENTER, first=True, after=4)
    para(tf, d, 11, MUTED, align=PP_ALIGN.CENTER, after=0)
    if i < 5:
        a = tb(s, x + sw + 0.01, 2.7, 0.22, 0.6, MSO_ANCHOR.MIDDLE)
        para(a, "›", 22, TEAL, True, align=PP_ALIGN.CENTER, first=True, after=0)
rect(s, 0.75, 4.15, 11.83, 0.8, RGBColor(0xB9, 0x1C, 0x1C), rounded=True)
tf = tb(s, 1.05, 4.32, 11.3, 0.5, MSO_ANCHOR.MIDDLE)
para(tf, [("security_check() runs ", False, WHITE), ("before every tool", True, WHITE), (", and every rejection is logged to rejection_log.jsonl", False, WHITE)], 15, first=True, after=0)
rect(s, 0.75, 5.2, 11.83, 0.85, SOFT, rounded=True); rect(s, 0.75, 5.2, 0.07, 0.85, TEAL)
tf = tb(s, 1.05, 5.38, 11.2, 0.5, MSO_ANCHOR.MIDDLE)
para(tf, "Known faults short-circuit at History: instant, no LLM, no network. Only new faults reach scrape and synthesis.", 14.5, RGBColor(0x47, 0x55, 0x69), first=True, after=0)
footer(s, 5)

# ---------- 6 THREE WAYS IN ----------
s = slide(); header(s, "Three ways in", "Upload a log, type a code, or batch a whole CSV")
img_fit(s, SHOT("05_batch.png"), 0.75, 2.1, 7.4, 4.4)
tf = tb(s, 8.5, 2.25, 4.1, 4.2)
bullet(tf, "A real IQAN system log, semicolon-delimited with a BOM, parsed cleanly", size=15)
bullet(tf, "De-duplicated to nine unique faults across the log", size=15)
bullet(tf, "One click: known faults instant, new ones scraped and reasoned", size=15)
rect(s, 8.5, 5.4, 4.08, 1.0, SOFT, rounded=True); rect(s, 8.5, 5.4, 0.06, 1.0, TEAL)
tf = tb(s, 8.72, 5.55, 3.75, 0.75, MSO_ANCHOR.MIDDLE)
para(tf, "Routine entries such as System started are filtered out automatically.", 12.5, RGBColor(0x47, 0x55, 0x69), first=True, after=0, line=1.2)
footer(s, 6)

# ---------- 7 CITED ANSWER ----------
s = slide(); header(s, "A cited answer", "Reasoned across multiple threads, never a single page")
rect(s, 0.75, 2.1, 6.3, 3.7, WHITE, LINE, rounded=True)
rect(s, 0.75, 2.1, 6.3, 0.09, TEAL)
tf = tb(s, 1.05, 2.35, 5.7, 0.5)
para(tf, [("XC44, No contact", True, NAVY)], 19, first=True, after=0)
rect(s, 5.45, 2.4, 1.45, 0.42, OKBG, rounded=True)
tf = tb(s, 5.45, 2.46, 1.45, 0.32, MSO_ANCHOR.MIDDLE)
para(tf, "High confidence", 10.5, OK, True, align=PP_ALIGN.CENTER, first=True, after=0)
tf = tb(s, 1.05, 3.0, 5.8, 2.6)
para(tf, [("Root cause: ", True, INK), ("the expansion module lost CAN communication, usually termination or wiring [1].", False, INK)], 14, first=True, after=8, line=1.2)
para(tf, "1.  Verify 120 ohm termination at both ends [1]", 14, INK, after=6)
para(tf, "2.  Inspect CAN_H and CAN_L for opens or shorts [1][2]", 14, INK, after=8)
para(tf, "Sources: No contact and critical CAN error, CAN bus-off termination", 12, MUTED, after=0)
tf = tb(s, 7.4, 2.35, 5.2, 3.4)
bullet(tf, "Inline citations [1][2] on every step, plus all source links", size=16)
bullet(tf, "Weighs higher-voted, answered community comments", size=16)
bullet(tf, "Output is parsed only as JSON. The model is never wired to a file or command tool", size=16)
footer(s, 7)

# ---------- 8 MAP ----------
s = slide(); header(s, "The knowledge map", "Related faults connect through shared forum threads")
tf = tb(s, 0.75, 2.4, 4.0, 3.6)
bullet(tf, "Each unique error is one node, sized by occurrence", size=16)
bullet(tf, "Linked to its fix and the forum threads it cited", size=16)
bullet(tf, "Shared threads converge, so related faults visibly connect", size=16)
img_fit(s, SHOT("06_map.png"), 5.0, 2.1, 7.6, 4.5)
footer(s, 8)

# ---------- 9 HARNESS ----------
s = slide(); header(s, "The harness", "Sandboxed tools, drop-in skills, editable config")
hcards = [("Sandboxed tools", "read, write, create, list, transform_markdown (md to HTML), and run_command. Every path passes through _safe_path()."),
          ("Drop-in skills", "Any skills/*.py with a SKILL dict is auto-discovered. Add one and it appears, with no harness edit."),
          ("Editable config", "mcp.json sets the LLM provider, forums, and scraping limits, with no code change.")]
for i, (t, d) in enumerate(hcards):
    card(s, 0.75 + i * 4.01, 2.15, 3.83, 2.5, t, d)
rect(s, 0.75, 5.05, 11.83, 0.95, SOFT, rounded=True); rect(s, 0.75, 5.05, 0.07, 0.95, TEAL)
tf = tb(s, 1.05, 5.25, 11.2, 0.6, MSO_ANCHOR.MIDDLE)
para(tf, "The loaded extension is extract_log (problem log to text and codes), alongside lookup_history, save_solution, phase1_crawl, and phase2_scrape.", 14, RGBColor(0x47, 0x55, 0x69), first=True, after=0)
footer(s, 9)

# ---------- 10 SECURITY LAYERS ----------
s = slide(); header(s, "Security, defense in depth", "A wide attack surface, guarded in four layers")
layers = [("1", "security_check(tool, args)", "Pattern and command screen before every tool call"),
          ("2", "_safe_path()", "Authoritative sandbox using resolve() and is_relative_to"),
          ("3", "check_url()", "Scheme check and domain allowlist before any HTTP"),
          ("4", "rejection_log.jsonl", "Every block recorded and shown in Settings")]
for i, (n, t, d) in enumerate(layers):
    x = 0.75 + (i % 2) * 6.04; y = 2.1 + (i // 2) * 1.2
    rect(s, x, y, 5.79, 1.05, WHITE, LINE, rounded=True)
    tf = tb(s, x + 0.25, y + 0.12, 0.6, 0.8, MSO_ANCHOR.MIDDLE)
    para(tf, n, 26, TEAL, True, first=True, after=0)
    tf = tb(s, x + 0.95, y + 0.16, 4.7, 0.8, MSO_ANCHOR.MIDDLE)
    para(tf, t, 15, NAVY, True, first=True, after=2)
    para(tf, d, 12, MUTED, after=0, line=1.1)
rect(s, 0.75, 4.65, 11.83, 1.35, NAVY, rounded=True)
tf = tb(s, 1.05, 4.9, 11.2, 0.95)
para(tf, [('security_check("run_command", {"command":"rm -rf /"})   ', False, RGBColor(0xE2, 0xE8, 0xF0)), ("returns (False, 'Blocked command: rm')", False, MUTED2)], 13, first=True, after=6)
para(tf, [('_safe_path("../../etc/passwd")   ', False, RGBColor(0xE2, 0xE8, 0xF0)), ("raises SandboxError", False, MUTED2)], 13, after=0)
footer(s, 10)

# ---------- 11 SECURITY VERDICTS ----------
s = slide(); header(s, "Security, per-tool verdicts", "Each risk is blocked, mitigated, or honestly accepted")
data = [("Tool or surface", "Risk", "Verdict"),
        ("File operations", "path traversal out of the sandbox", "Blocked"),
        ("run_command", "shell injection from tool output", "Mitigated"),
        ("transform_markdown", "HTML or script injection", "Mitigated"),
        ("Outbound HTTP", "SSRF and scraping abuse", "Mitigated"),
        ("Prompt injection", "scraped or file text as instructions", "By design"),
        ("Secrets", "API key leakage", "Mitigated"),
        ("Skills folder", "arbitrary code at import time", "Accepted")]
gt = s.shapes.add_table(len(data), 3, Inches(0.75), Inches(2.05), Inches(11.83), Inches(3.7)).table
gt.columns[0].width = Inches(3.2); gt.columns[1].width = Inches(5.9); gt.columns[2].width = Inches(2.73)
for ri, rowv in enumerate(data):
    for ci, val in enumerate(rowv):
        cell = gt.cell(ri, ci)
        cell.margin_left = Inches(0.12); cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = val
        f = r.font; f.size = Pt(14); f.name = "Segoe UI"
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY; f.color.rgb = WHITE; f.bold = True
        else:
            f.color.rgb = INK
            if ci == 0:
                f.bold = True
            if ci == 2:
                acc = WARN if val == "Accepted" else OK
                accbg = WARNBG if val == "Accepted" else OKBG
                cell.fill.solid(); cell.fill.fore_color.rgb = accbg; f.color.rgb = acc; f.bold = True
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else SOFT
tf = tb(s, 0.75, 5.95, 11.83, 0.5)
para(tf, "Honest acceptance of a risk that is understood, rather than pretending it is not there.", 13.5, MUTED, first=True, after=0)
footer(s, 11)

# ---------- 12 DELIVERABLES ----------
s = slide(); header(s, "How it maps to the assignment", "All four deliverables are in the repository")
dels = [("The harness", "Working code. The app runs headless with zero exceptions"),
        ("A loaded extension", "Five drop-in skills, led by extract_log"),
        ("README and security", "Per-tool risks marked blocked, mitigated, or accepted, with diagrams"),
        ("Presentation", "This deck, with screenshots from the live app")]
for i, (t, d) in enumerate(dels):
    x = 0.75 + (i % 2) * 6.04; y = 2.1 + (i // 2) * 1.45
    rect(s, x, y, 5.79, 1.3, WHITE, LINE, rounded=True)
    tf = tb(s, x + 0.28, y + 0.2, 5.2, 1.0)
    para(tf, [("✓  ", True, TEAL), (t, True, NAVY)], 17, first=True, after=4)
    para(tf, d, 13, MUTED, after=0, line=1.15)
tf = tb(s, 0.75, 5.2, 12, 0.5)
para(tf, [("Offline (Ollama)     Safety interlock and log     Evaluation and retry loop     Context caps     Staged multi-agent roles", True, TEAL)], 13.5, first=True, after=0)
footer(s, 12)

# ---------- 13 HONESTY ----------
s = slide(); header(s, "Honesty notes", "Stated plainly, because the brief rewards it")
tf = tb(s, 0.75, 2.15, 11.8, 4.3)
bullet(tf, "The real extension mechanism is the drop-in skills loader.", lead="mcp.json is config, not a literal MCP server. ", size=16.5, after=12)
bullet(tf, "They are trusted first-party code; adding one is equivalent to editing the app.", lead="Skills are not sandboxed. ", size=16.5, after=12)
bullet(tf, "It is fully capable, an accepted convenience.", lead="python is allow-listed for run_command. ", size=16.5, after=12)
bullet(tf, "The crawler rate-limits, identifies itself, and reads only public threads.", lead="No robots.txt parsing. ", size=16.5, after=12)
bullet(tf, "Streamlit Cloud resets on reboot; persistence would need a database.", lead="Hosted history is ephemeral. ", size=16.5, after=0)
footer(s, 13)

# ---------- 14 DEPLOYMENT ----------
s = slide(); header(s, "Deployment", "Runs in two modes")
rect(s, 0.75, 2.2, 5.8, 2.2, NAVY, rounded=True)
tf = tb(s, 1.05, 2.45, 5.3, 1.8)
para(tf, "# Offline, local, no key", 13, MUTED2, first=True, after=8)
para(tf, "ollama serve  &&  ollama pull phi3", 13.5, RGBColor(0xE2, 0xE8, 0xF0), after=6)
para(tf, "streamlit run app.py", 13.5, RGBColor(0xE2, 0xE8, 0xF0), after=4)
para(tf, "# toggle Online off", 12.5, MUTED2, after=0)
rect(s, 6.78, 2.2, 5.8, 2.2, NAVY, rounded=True)
tf = tb(s, 7.08, 2.45, 5.3, 1.8)
para(tf, "# Online, hosted LLM", 13, MUTED2, first=True, after=8)
para(tf, "export OPENAI_API_KEY=sk-...", 13.5, RGBColor(0xE2, 0xE8, 0xF0), after=6)
para(tf, "streamlit run app.py", 13.5, RGBColor(0xE2, 0xE8, 0xF0), after=4)
para(tf, "# toggle Online on", 12.5, MUTED2, after=0)
tf = tb(s, 0.75, 4.7, 11.8, 0.6)
para(tf, "The hosted instance runs live at rootiq.streamlit.app, with the API key held in Streamlit secrets.", 15, INK, first=True, after=0)
links_line(s, 0.75, 5.45, color=INK)
footer(s, 14)

# ---------- 15 CLOSE ----------
s = slide(); rect(s, 0, 0, 13.333, 7.5, NAVY); rect(s, 0, 0, 0.22, 7.5, TEAL)
tf = tb(s, 0.95, 2.4, 11, 1.4)
para(tf, "Thank you", 56, WHITE, True, first=True, after=0)
tf = tb(s, 0.97, 3.75, 11, 0.6)
para(tf, "Secure, dual-mode PLC error intelligence", 22, RGBColor(0xE2, 0xE8, 0xF0), first=True, after=0)
tf = tb(s, 0.97, 4.7, 11, 0.5)
para(tf, [("Victor Otu Hayford", True, WHITE), (", Industrial AI, 2026", True, TEALL)], 19, first=True, after=0)
links_line(s, 0.97, 5.4)

OUT = "RootIQ-Presentation.pptx"
prs.save(OUT)
print(f"Saved {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
